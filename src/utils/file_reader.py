# src/utils/file_reader.py
import io
import base64
from pathlib import Path
from typing import Dict

import pandas as pd
from docx import Document
import pdfplumber
from pptx import Presentation
from PIL import Image

# CSV
def read_csv_file(bytes_: bytes) -> str:
    text = bytes_.decode("utf-8", errors="ignore")
    # limitamos tamanho para não explodir o prompt
    return text[:5000]

# Excel
def read_excel_file(bytes_: bytes) -> str:
    df = pd.read_excel(io.BytesIO(bytes_))
    return df.head(10).to_markdown(index=False)

# TXT / MD / RTF
def read_txt_file(bytes_: bytes) -> str:
    return bytes_.decode("utf-8", errors="ignore")

# DOCX
def read_docx_file(bytes_: bytes) -> str:
    doc = Document(io.BytesIO(bytes_))
    return "\n".join([para.text for para in doc.paragraphs])

# PDF
def read_pdf_file(bytes_: bytes) -> str:
    text = ""
    with pdfplumber.open(io.BytesIO(bytes_)) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

# PPTX
def read_pptx_file(bytes_: bytes) -> str:
    prs = Presentation(io.BytesIO(bytes_))
    slides_text = []
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        slides_text.append("\n".join(texts))
    return "\n\n".join(slides_text)

# Imagem – devolve string base64 (não gera texto real)
def read_image_file(bytes_: bytes) -> str:
    return base64.b64encode(bytes_).decode("utf-8")
